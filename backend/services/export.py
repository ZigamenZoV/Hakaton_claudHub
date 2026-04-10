from __future__ import annotations

import io
import json
import re

from services.mws_client import chat_complete


async def _llm_structure(content: str, title: str) -> list[dict]:
    prompt = f"""Convert this content into a presentation structure. Output ONLY a JSON array of slides, each with:
{{"title": "slide title", "bullets": ["point 1", "point 2"], "notes": "speaker notes"}}

Title: {title}
Content:
{content[:3000]}

Output 5-7 slides maximum. JSON only, no other text."""

    raw = await chat_complete([{"role": "user", "content": prompt}])
    raw = re.sub(r"```json|```", "", raw).strip()
    try:
        return json.loads(raw)
    except Exception:
        return [{"title": title, "bullets": [content[:200]], "notes": ""}]


async def to_pptx(content: str, title: str = "GPTHub Report") -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    slides_data = await _llm_structure(content, title)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Сгенерировано GPTHub"

    for sd in slides_data:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = sd.get("title", "")
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        for i, bullet in enumerate(sd.get("bullets", [])):
            if i == 0:
                tf.paragraphs[0].text = bullet
            else:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
        if sd.get("notes"):
            slide.notes_slide.notes_text_frame.text = sd["notes"]

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


async def to_pdf(content: str, title: str = "GPTHub Report") -> bytes:
    import markdown as md_lib
    from weasyprint import HTML

    html_body = md_lib.markdown(content, extensions=["tables", "fenced_code"])
    html = f"""<!DOCTYPE html>
<html>
<head>
<meta charset="utf-8">
<title>{title}</title>
<style>
  body {{ font-family: 'Helvetica Neue', Arial, sans-serif; margin: 2cm; color: #1a1a1a; }}
  h1 {{ color: #e31837; border-bottom: 2px solid #e31837; padding-bottom: 8px; }}
  h2 {{ color: #333; margin-top: 1.5em; }}
  code {{ background: #f5f5f5; padding: 2px 6px; border-radius: 3px; font-size: 0.9em; }}
  pre {{ background: #f5f5f5; padding: 1em; border-radius: 6px; overflow-x: auto; }}
  table {{ border-collapse: collapse; width: 100%; }}
  th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
  th {{ background: #e31837; color: white; }}
</style>
</head>
<body>
<h1>{title}</h1>
{html_body}
</body>
</html>"""

    return HTML(string=html).write_pdf()